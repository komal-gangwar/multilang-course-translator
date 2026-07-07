"""
============================================================
  Document Processing Utilities
  File: app/utils/document_processor.py
============================================================
"""

import io
import re
import logging
from pathlib import Path
from typing import List, Tuple

logger = logging.getLogger(__name__)


# ── Text cleaning ────────────────────────────────────────────────────────────

def clean_text(text: str) -> str:
    """Remove excessive whitespace, fix encoding artefacts."""
    text = re.sub(r'\r\n', '\n', text)
    text = re.sub(r'\r', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    return text.strip()


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """
    Split text into overlapping chunks for RAG indexing.

    Args:
        text: Full document text
        chunk_size: Target chunk size in characters
        overlap: Overlap between consecutive chunks

    Returns:
        List of text chunks
    """
    paragraphs = re.split(r'\n\s*\n', text)
    chunks = []
    current = ""

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        if len(current) + len(para) <= chunk_size:
            current = (current + "\n\n" + para).strip() if current else para
        else:
            if current:
                chunks.append(current)
            # If paragraph itself is too long, split by sentences
            if len(para) > chunk_size:
                sentences = re.split(r'(?<=[.!?])\s+', para)
                sub = ""
                for sent in sentences:
                    if len(sub) + len(sent) <= chunk_size:
                        sub = (sub + " " + sent).strip() if sub else sent
                    else:
                        if sub:
                            chunks.append(sub)
                        sub = sent
                if sub:
                    current = sub
                else:
                    current = ""
            else:
                current = para

    if current:
        chunks.append(current)

    # Add overlap
    overlapped = []
    for i, chunk in enumerate(chunks):
        if i > 0 and overlap > 0:
            prev_words = chunks[i - 1].split()[-overlap // 5:]
            prefix = " ".join(prev_words)
            chunk = (prefix + " " + chunk).strip()
        overlapped.append(chunk)

    return overlapped


# ── File type detection ──────────────────────────────────────────────────────

ALLOWED_EXTENSIONS = {"pdf", "pptx", "ppt", "docx", "doc", "txt", "md"}


def allowed_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def get_file_type(filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "unknown"
    return ext


# ── PDF extraction ───────────────────────────────────────────────────────────

def extract_pdf(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extract text from a PDF file.

    Returns:
        (full_text, page_count) tuple
    """
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return clean_text("\n\n".join(pages)), len(pages)
    except ImportError:
        logger.error("PyMuPDF not installed. Run: pip install PyMuPDF")
        return "", 0
    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        return "", 0


# ── PPTX extraction ──────────────────────────────────────────────────────────

def extract_pptx(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extract text from a PowerPoint file.

    Returns:
        (full_text, slide_count) tuple
    """
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            parts.append(line)
            slides_text.append("\n".join(parts))
        return clean_text("\n\n".join(slides_text)), len(prs.slides)
    except ImportError:
        logger.error("python-pptx not installed. Run: pip install python-pptx")
        return "", 0
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return "", 0


# ── DOCX extraction ──────────────────────────────────────────────────────────

def extract_docx(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extract text from a Word document.

    Returns:
        (full_text, paragraph_count) tuple
    """
    try:
        import docx
        doc = docx.Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return clean_text("\n\n".join(paragraphs)), len(paragraphs)
    except ImportError:
        logger.error("python-docx not installed. Run: pip install python-docx")
        return "", 0
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        return "", 0


# ── Plain text extraction ────────────────────────────────────────────────────

def extract_text_file(file_bytes: bytes) -> Tuple[str, int]:
    """Decode and clean a plain text / Markdown file."""
    try:
        import chardet
        detected = chardet.detect(file_bytes)
        encoding = detected.get("encoding", "utf-8") or "utf-8"
        text = file_bytes.decode(encoding, errors="replace")
        lines = [l for l in text.splitlines() if l.strip()]
        return clean_text(text), len(lines)
    except Exception as e:
        logger.error(f"Text file extraction failed: {e}")
        return "", 0


# ── Dispatcher ───────────────────────────────────────────────────────────────

def extract_text(file_bytes: bytes, filename: str) -> Tuple[str, int, str]:
    """
    Dispatch extraction based on file extension.

    Returns:
        (text, page_or_slide_count, file_type) tuple
    """
    ext = get_file_type(filename)

    if ext == "pdf":
        text, count = extract_pdf(file_bytes)
        return text, count, "pdf"
    elif ext in ("pptx", "ppt"):
        text, count = extract_pptx(file_bytes)
        return text, count, "presentation"
    elif ext in ("docx", "doc"):
        text, count = extract_docx(file_bytes)
        return text, count, "document"
    elif ext in ("txt", "md"):
        text, count = extract_text_file(file_bytes)
        return text, count, "text"
    else:
        return "", 0, "unknown"


# ── Word count & preview ─────────────────────────────────────────────────────

def word_count(text: str) -> int:
    return len(text.split())


def preview(text: str, chars: int = 400) -> str:
    """Return a short preview of the text."""
    if len(text) <= chars:
        return text
    return text[:chars].rsplit(" ", 1)[0] + "…"
